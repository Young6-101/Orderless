import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from backend.file_processor.image_handler import ImageDescriber


def extract_text_from_pptx(file_path: str, describe_images: bool = True) -> str:
    """
    Extract text + image descriptions from a .pptx file.

    For each slide:
      - Extracts all text from text frames (titles, body, tables, etc.)
      - Optionally uses ImageDescriber to generate descriptions for embedded images.

    Args:
        file_path: Path to the .pptx file.
        describe_images: If True, use vision model to describe embedded images.

    Returns:
        str: Combined text content from the entire presentation.
    """
    path = Path(file_path)
    if not path.exists():
        print(f"--- [ERROR] PPTX file not found: {file_path} ---")
        return ""

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"--- [ERROR] Failed to open PPTX {file_path}: {e} ---")
        return ""

    describer = ImageDescriber() if describe_images else None
    all_parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = [f"[Slide {slide_idx}]"]

        for shape in slide.shapes:
            # --- Text extraction ---
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)

            # --- Table extraction ---
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells
                    )
                    if row_text.strip("| "):
                        slide_parts.append(row_text)

            # --- Image extraction & description ---
            if (
                describe_images
                and describer
                and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    # Save to a temp file for the describer
                    ext = image.content_type.split("/")[-1]  # e.g. "png"
                    tmp_img_path = Path(f"temp_slide{slide_idx}_img.{ext}")
                    tmp_img_path.write_bytes(image_bytes)

                    desc = describer.describe_image(
                        str(tmp_img_path),
                        prompt=(
                            "Describe this presentation slide image. "
                            "Focus on diagrams, charts, or visual content."
                        ),
                    )
                    slide_parts.append(f"[Image on slide {slide_idx}: {desc}]")

                    # Cleanup temp image
                    tmp_img_path.unlink(missing_ok=True)
                except Exception as e:
                    print(
                        f"--- [WARN] Could not describe image on slide {slide_idx}: {e} ---"
                    )

        slide_text = "\n".join(slide_parts)
        if slide_text.strip():
            all_parts.append(slide_text)

    result = "\n\n".join(all_parts)
    clean_result = " ".join(result.split())
    print(f"--- [LOG] Extracted {len(prs.slides)} slides from {path.name} ---")
    return clean_result
